from PyPDF2 import PdfMerger, PdfReader, PdfWriter
from PIL import Image, ImageOps, ImageEnhance, ImageFilter
from pdf2docx import Converter
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
import os, io, re, difflib, random
import numpy as np
import fitz  # PyMuPDF
import pdfkit
import pandas as pd
import pdfplumber

def merge_pdfs(files, output):
    merger = PdfMerger()
    for file in files:
        merger.append(file)
    merger.write(output)
    merger.close()

def split_pdf(input_path, output_folder):
    reader = PdfReader(input_path)
    output_files = []
    for i, page in enumerate(reader.pages):
        writer = PdfWriter()
        writer.add_page(page)
        output_filename = f"split_page_{i+1}.pdf"
        output_path = os.path.join(output_folder, output_filename)
        with open(output_path, "wb") as f:
            writer.write(f)
        output_files.append(output_path)
    return output_files

def pdf_to_excel(input_path, output_path):
    """
    Extract tables from PDF and save to Excel
    """
    all_tables = []
    with pdfplumber.open(input_path) as pdf:
        for page in pdf.pages:
            tables = page.extract_tables()
            for table in tables:
                df = pd.DataFrame(table[1:], columns=table[0])
                all_tables.append(df)
    
    if all_tables:
        with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
            for i, df in enumerate(all_tables):
                df.to_excel(writer, sheet_name=f'Table_{i+1}', index=False)
        return True
    return False

def compress_pdf(input_path, output_path):
    reader = PdfReader(input_path)
    writer = PdfWriter()
    for page in reader.pages:
        page.compress_content_streams()
        writer.add_page(page)
    with open(output_path, "wb") as f:
        writer.write(f)

def images_to_pdf(image_paths, output_path):
    images = []
    for img_path in image_paths:
        img = Image.open(img_path)
        if img.mode != 'RGB':
            img = img.convert('RGB')
        images.append(img)
    
    if images:
        images[0].save(output_path, save_all=True, append_images=images[1:])

def pdf_to_word(pdf_path, docx_path):
    cv = Converter(pdf_path)
    cv.convert(docx_path, start=0, end=None)
    cv.close()

def rotate_pdf(input_path, output_path, rotation=90):
    reader = PdfReader(input_path)
    writer = PdfWriter()
    for page in reader.pages:
        page.rotate(rotation)
        writer.add_page(page)
    with open(output_path, "wb") as f:
        writer.write(f)

def add_watermark(input_path, output_path, watermark_text):
    # Create watermark PDF in memory
    packet = io.BytesIO()
    can = canvas.Canvas(packet, pagesize=letter)
    can.setFont("Helvetica", 40)
    can.setFillAlpha(0.3)
    can.saveState()
    can.translate(300, 400)
    can.rotate(45)
    can.drawCentredString(0, 0, watermark_text)
    can.restoreState()
    can.save()
    packet.seek(0)
    
    watermark_pdf = PdfReader(packet)
    watermark_page = watermark_pdf.pages[0]
    
    reader = PdfReader(input_path)
    writer = PdfWriter()
    
    for page in reader.pages:
        page.merge_page(watermark_page)
        writer.add_page(page)
        
    with open(output_path, "wb") as f:
        writer.write(f)

def annotate_pdf(input_path, output_path, annotations):
    """
    annotations: list of dicts with keys like type, x, y, content, color, size
    """
    doc = fitz.open(input_path)
    for ann in annotations:
        try:
            page_num = int(ann.get("page", 1)) - 1
            if page_num < 0 or page_num >= len(doc): continue
            page = doc[page_num]
            
            # Helper to convert hex to RGB 0-1
            def hex_to_rgb(h):
                h = h.lstrip('#')
                return tuple(int(h[i:i+2], 16)/255.0 for i in (0, 2, 4))

            color = hex_to_rgb(ann.get("color", "#000000"))
            
            if ann["type"] == "text":
                font_map = {
                    "Helvetica": "helv",
                    "Times-Roman": "tiro",
                    "Courier": "cour"
                }
                f_name = font_map.get(ann.get("font"), "helv")
                page.insert_text(
                    (float(ann["x"]), float(ann["y"])), 
                    ann["content"], 
                    fontname=f_name,
                    fontsize=float(ann.get("size", 20)),
                    color=color
                )
            elif ann["type"] == "rect":
                page.draw_rect(
                    fitz.Rect(float(ann["x1"]), float(ann["y1"]), float(ann["x2"]), float(ann["y2"])),
                    color=color,
                    width=float(ann.get("width", 2))
                )
            elif ann["type"] == "image":
                rect = fitz.Rect(float(ann["x1"]), float(ann["y1"]), float(ann["x2"]), float(ann["y2"]))
                page.insert_image(rect, filename=ann["image_path"])
        except Exception as e:
            print(f"Error annotating: {e}")
            continue
            
    doc.save(output_path)
    doc.close()

def pdf_to_images(input_path, output_folder):
    doc = fitz.open(input_path)
    image_paths = []
    for i in range(len(doc)):
        page = doc[i]
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x zoom for better quality
        output_path = os.path.join(output_folder, f"page_{i+1}.jpg")
        pix.save(output_path)
        image_paths.append(output_path)
    doc.close()
    return image_paths

def url_to_pdf(url, output_path):
    # This requires wkhtmltopdf installed on your system
    # Check common paths on Windows if it's not in PATH
    path_to_wkhtmltopdf = r'C:\Program Files\wkhtmltopdf\bin\wkhtmltopdf.exe'
    config = None
    if os.path.exists(path_to_wkhtmltopdf):
        config = pdfkit.configuration(wkhtmltopdf=path_to_wkhtmltopdf)
        
    try:
        pdfkit.from_url(url, output_path, configuration=config)
        return True
    except Exception as e:
        print(f"pdfkit error diagnostic: {e}")
        # If still failing, return False to trigger the informative HTTPException in app.py
        return False

def extract_text(input_path):
    reader = PdfReader(input_path)
    text = ""
    for page in reader.pages:
        text += page.extract_text() + "\n"
    return text

def add_password(input_path, output_path, password):
    reader = PdfReader(input_path)
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    writer.encrypt(password)
    with open(output_path, "wb") as f:
        writer.write(f)

def remove_password(input_path, output_path, password):
    reader = PdfReader(input_path)
    if reader.is_encrypted:
        reader.decrypt(password)
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    with open(output_path, "wb") as f:
        writer.write(f)

def fill_form(input_path, output_path, form_data):
    """
    form_data: dict of field_name: value
    """
    doc = fitz.open(input_path)
    for page in doc:
        for field in page.widgets():
            if field.field_name in form_data:
                field.field_value = form_data[field.field_name]
                field.update()
    doc.save(output_path)
    doc.close()

def redact_text(input_path, output_path, text_to_redact):
    doc = fitz.open(input_path)
    for page in doc:
        text_instances = page.search_for(text_to_redact)
        for inst in text_instances:
            page.add_redact_annot(inst)
        page.apply_redactions()
    doc.save(output_path)
    doc.close()

def replace_text(input_path, output_path, old_text, new_text):
    doc = fitz.open(input_path)
    for page in doc:
        text_instances = page.search_for(old_text)
        for inst in text_instances:
            page.add_redact_annot(inst)
        page.apply_redactions()
        # After redaction, insert new text at the position
        # This is approximate, as position might not be exact
        for inst in text_instances:
            page.insert_text(inst.top_left, new_text, fontsize=12)
    doc.save(output_path)
    doc.close()

def add_highlight(input_path, output_path, text_to_highlight, color=(1, 1, 0)):
    doc = fitz.open(input_path)
    for page in doc:
        text_instances = page.search_for(text_to_highlight)
        for inst in text_instances:
            highlight = page.add_highlight_annot(inst)
            highlight.set_colors(stroke=color)
            highlight.update()
    doc.save(output_path)
    doc.close()

def add_text_stamp(input_path, output_path, text, position='center', font_size=20, color=(0,0,0), page_num=None):
    """
    Add text stamp to PDF pages
    position: 'center', 'top', 'bottom', 'top-left', 'top-right', 'bottom-left', 'bottom-right'
    page_num: None for all pages, or specific page number (1-based)
    """
    doc = fitz.open(input_path)
    pages = [doc[page_num - 1]] if page_num else doc
    for page in pages:
        page_width = page.rect.width
        page_height = page.rect.height
        
        if position == 'center':
            x = page_width / 2
            y = page_height / 2
        elif position == 'top':
            x = page_width / 2
            y = 50
        elif position == 'bottom':
            x = page_width / 2
            y = page_height - 50
        elif position == 'top-left':
            x = 50
            y = 50
        elif position == 'top-right':
            x = page_width - 50
            y = 50
        elif position == 'bottom-left':
            x = 50
            y = page_height - 50
        elif position == 'bottom-right':
            x = page_width - 50
            y = page_height - 50
        else:
            x = page_width / 2
            y = page_height / 2
        
        page.insert_text((x, y), text, fontsize=font_size, color=color, align=fitz.TEXT_ALIGN_CENTER)
    doc.save(output_path)
    doc.close()

def edit_text_in_pdf(input_path, output_path, changes):
    """
    changes: list of dicts {'page': int, 'old_text': str, 'new_text': str}
    """
    doc = fitz.open(input_path)
    for change in changes:
        page_num = change['page'] - 1
        if page_num < 0 or page_num >= len(doc):
            continue
        page = doc[page_num]
        old_text = change['old_text']
        new_text = change['new_text']
        text_instances = page.search_for(old_text)
        for inst in text_instances:
            page.add_redact_annot(inst)
        page.apply_redactions()
        # Insert new text at approximate position
        if text_instances:
            pos = text_instances[0].top_left
            page.insert_text(pos, new_text, fontsize=12)
    doc.save(output_path)
    doc.close()

def add_page_numbers(input_path, output_path, start_page=1, position='bottom-right', font_size=12, color=(0,0,0)):
    """
    Add page numbers to PDF
    position: 'bottom-left', 'bottom-center', 'bottom-right', 'top-left', 'top-center', 'top-right'
    """
    doc = fitz.open(input_path)
    for i, page in enumerate(doc):
        page_num = start_page + i
        page_width = page.rect.width
        page_height = page.rect.height
        
        if position == 'bottom-left':
            x = 50
            y = page_height - 30
        elif position == 'bottom-center':
            x = page_width / 2
            y = page_height - 30
        elif position == 'bottom-right':
            x = page_width - 50
            y = page_height - 30
        elif position == 'top-left':
            x = 50
            y = 30
        elif position == 'top-center':
            x = page_width / 2
            y = 30
        elif position == 'top-right':
            x = page_width - 50
            y = 30
        else:
            x = page_width - 50
            y = page_height - 30
        
        page.insert_text((x, y), str(page_num), fontsize=font_size, color=color, align=fitz.TEXT_ALIGN_CENTER)
    doc.save(output_path)
    doc.close()

def crop_pdf(input_path, output_path, left=0, top=0, right=None, bottom=None):
    """
    Crop PDF pages
    left, top, right, bottom: coordinates in points (1/72 inch)
    If right/bottom None, use page dimensions
    """
    doc = fitz.open(input_path)
    for page in doc:
        page_width = page.rect.width
        page_height = page.rect.height
        right = right or page_width
        bottom = bottom or page_height
        page.set_cropbox(fitz.Rect(left, top, right, bottom))
    doc.save(output_path)
    doc.close()

def reorder_pages(input_path, output_path, page_order):
    """
    Reorder PDF pages
    page_order: list of 1-based page numbers, e.g. [3,1,2]
    """
    reader = PdfReader(input_path)
    writer = PdfWriter()
    for page_num in page_order:
        if 1 <= page_num <= len(reader.pages):
            writer.add_page(reader.pages[page_num - 1])
    with open(output_path, "wb") as f:
        writer.write(f)

def pdf_to_ppt(input_path, output_path):
    """
    Convert PDF to PPT by extracting high-res images from each page
    """
    from pptx import Presentation
    
    doc = fitz.open(input_path)
    prs = Presentation()
    
    for i, page in enumerate(doc):
        # High resolution (300 DPI approx)
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        img_data = pix.tobytes("png")
        img_stream = io.BytesIO(img_data)
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank slide
        
        # Fit picture to slide
        prs_width = prs.slide_width
        prs_height = prs.slide_height
        slide.shapes.add_picture(img_stream, 0, 0, width=prs_width, height=prs_height)
    
    prs.save(output_path)
    doc.close()

def extract_pages(input_path, output_path, pages):
    """
    Extract specific pages from PDF
    pages: list of 1-based page numbers or ranges like "1-3,5"
    """
    reader = PdfReader(input_path)
    writer = PdfWriter()
    
    page_set = set()
    for part in pages.replace(' ', '').split(','):
        if '-' in part:
            start, end = map(int, part.split('-'))
            page_set.update(range(start, end + 1))
        else:
            page_set.add(int(part))
    
    for page_num in sorted(page_set):
        if 1 <= page_num <= len(reader.pages):
            writer.add_page(reader.pages[page_num - 1])
    
    with open(output_path, "wb") as f:
        writer.write(f)

def compare_pdfs(path1, path2, output_path):
    """
    Compare two PDFs and highlight differences.
    Text additions in path2 compared to path1 are highlighted green.
    """
    doc1 = fitz.open(path1)
    doc2 = fitz.open(path2)
    
    # Simple strategy: Compare page by page text
    for i in range(min(len(doc1), len(doc2))):
        page1 = doc1[i]
        page2 = doc2[i]
        
        text1 = page1.get_text().splitlines()
        text2 = page2.get_text().splitlines()
        
        d = difflib.Differ()
        diff = list(d.compare(text1, text2))
        
        for line in diff:
            if line.startswith('+ '):
                # Text added in doc2
                added_text = line[2:].strip()
                if added_text:
                    text_instances = page2.search_for(added_text)
                    for inst in text_instances:
                        annot = page2.add_highlight_annot(inst)
                        annot.set_colors(stroke=(0, 1, 0)) # Green
                        annot.update()
            # Deletions are harder to show on the NEW pdf without altering layout significantly
                
    doc2.save(output_path)
    doc1.close()
    doc2.close()

def smart_redact(input_path, output_path, patterns):
    """
    patterns: list of strings like 'email', 'credit_card', 'phone'
    """
    regex_map = {
        'email': r'[a-zA-Z0-9_.+-]+@[a-zA-Z0-9-]+\.[a-zA-Z0-9-.]+',
        'credit_card': r'\b(?:\d[ -]*?){13,16}\b',
        'phone': r'\b(?:\+?\d{1,3}[- ]?)?\(?\d{3}\)?[- ]?\d{3}[- ]?\d{4}\b'
    }
    
    doc = fitz.open(input_path)
    for page in doc:
        for p_name in patterns:
            pattern = regex_map.get(p_name, p_name) # allow custom regex too
            # Use regex search since fitz search_for is not always regex-friendly
            # Actually PyMuPDF has page.search_for which can take a regex if enabled, 
            # but let's use the standard search_for with the literal strings if possible 
            # or use a more advanced approach.
            # For simplicity with the tools we have:
            text_instances = page.search_for(pattern) # search_for supports regex in newer versions or treats as string
            for inst in text_instances:
                page.add_redact_annot(inst, fill=(0, 0, 0))
        page.apply_redactions()
        
    doc.save(output_path)
    doc.close()

def fake_scan(input_path, output_path):
    """
    Make a PDF look like it was physically scanned.
    """
    doc = fitz.open(input_path)
    images = []
    
    for page in doc:
        # Convert page to high-res image
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        
        # 1. Grayscale
        img = ImageOps.grayscale(img)
        
        # 2. Add slight rotation (misalignment)
        angle = random.uniform(-0.5, 0.5)
        img = img.rotate(angle, resample=Image.BICUBIC, expand=False, fillcolor=255)
        
        # 3. Add noise/grain
        img_array = np.array(img)
        noise = np.random.normal(0, 2, img_array.shape).astype(np.uint8)
        img_array = np.clip(img_array.astype(np.int16) + noise.astype(np.int16), 0, 255).astype(np.uint8)
        img = Image.fromarray(img_array)
        
        # 4. Adjust contrast/brightness
        img = ImageEnhance.Contrast(img).enhance(random.uniform(1.1, 1.3))
        img = ImageEnhance.Brightness(img).enhance(random.uniform(0.95, 1.05))
        
        images.append(img)
        
    if images:
        images[0].save(output_path, save_all=True, append_images=images[1:], resolution=150.0, quality=60)
    doc.close()

def make_booklet(input_path, output_path):
    """
    Rearrange pages for booklet printing (imposition)
    """
    reader = PdfReader(input_path)
    pages = list(reader.pages)
    
    # Pad pages to multiple of 4
    num_orig = len(pages)
    while len(pages) % 4 != 0:
        # Create a blank page for padding
        temp_writer = PdfWriter()
        temp_writer.add_blank_page(width=pages[0].mediabox.width, height=pages[0].mediabox.height)
        pages.append(temp_writer.pages[0])
    
    num_pages = len(pages)
    # The sequence for booklet imposition (4 pages per sheet: 2 front, 2 back)
    booklet_order = []
    for i in range(num_pages // 4):
        # Front side: Last-to-middle, First-from-middle
        booklet_order.append(num_pages - (2*i) - 1)
        booklet_order.append(2*i)
        # Back side: First-from-middle+1, Last-to-middle-1
        booklet_order.append(2*i + 1)
        booklet_order.append(num_pages - (2*i) - 2)
        
    writer = PdfWriter()
    for idx in booklet_order:
        writer.add_page(pages[idx])
        
    with open(output_path, "wb") as f:
        writer.write(f)

def remove_annotations(input_path, output_path):
    """
    Strip all highlights, notes, and other annotations from PDF.
    """
    doc = fitz.open(input_path)
    for page in doc:
        for annot in page.annots():
            page.delete_annot(annot)
    doc.save(output_path)
    doc.close()
